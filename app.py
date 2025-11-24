import streamlit as st
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
import io
from datetime import datetime
import os

# --- ตั้งค่าหน้าเว็บ ---
st.set_page_config(page_title="ระบบสร้างรายงานสายตรวจทางน้ำ", layout="wide")
st.title("👮‍♂️ ระบบสร้างรายงานสายตรวจท่าเรือประจำวัน (ส.รน.4)")

# --- ตรวจสอบไฟล์จำเป็น ---
required_files = {
    "template": "Template.pptx",
    "background": "background.jpg",   # <--- ต้องมีรูปพื้นหลังเปล่า
    "font": "THSarabunNew.ttf"        # <--- ต้องมีไฟล์ฟอนต์
}

missing_files = [f for f in required_files.values() if not os.path.exists(f)]
if missing_files:
    st.error(f"❌ ไฟล์ไม่ครบ! กรุณาอัปโหลดไฟล์เหล่านี้ขึ้นระบบ: {', '.join(missing_files)}")
    st.stop()

# --- ฟังก์ชันแก้ไขข้อความใน PPT (เหมือนเดิม) ---
def replace_text_ppt(shape, search_str, replace_str):
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            full_text = "".join([run.text for run in paragraph.runs])
            if search_str in full_text:
                for run in paragraph.runs:
                     if search_str in run.text:
                        run.text = run.text.replace(search_str, replace_str)

# --- ฟังก์ชันสร้างรูปภาพ (สำหรับ JPG/PDF) ---
def generate_image_report(data, images, bg_path, font_path):
    # 1. เปิดรูปพื้นหลัง
    base_img = Image.open(bg_path).convert("RGB")
    draw = ImageDraw.Draw(base_img)
    
    # 2. โหลดฟอนต์ (ปรับขนาดตามความเหมาะสม)
    try:
        font_header = ImageFont.truetype(font_path, 50) # ฟอนต์หัวข้อใหญ่
        font_text = ImageFont.truetype(font_path, 28)   # ฟอนต์เนื้อหา
    except:
        st.error("โหลดฟอนต์ไม่ได้ เช็คชื่อไฟล์ฟอนต์ให้ถูกต้อง")
        return None

    # 3. กำหนดพิกัดข้อความ (X, Y) - **ต้องปรับแก้ตัวเลขตรงนี้ให้ตรงกับช่องว่างในรูปของคุณ**
    # (ตัวเลขสมมติ: แกน X แนวนอน, แกน Y แนวตั้ง)
    text_color = (0, 0, 0) # สีดำ

    # --- ส่วนที่แก้ไข: ปรับพิกัดข้อความ (แก้ Syntax Error แล้ว) ---

    # 1. เขียนหัวข้อเดือน (มุมขวาบน)
    draw.text((400, 90), data["{{HEADER_MONTH}}"], font=font_header, fill=(255, 255, 0)) 

    # 2. ตั้งค่าระยะห่าง
    # start_x: ขยับไปขวา 1150 เพื่อหลบหัวข้อ
    start_x = 980  
    
    # start_y: บรรทัดแรกเริ่มที่ความสูงนี้
    start_y = 200   
    
    # gap: ระยะห่างระหว่างบรรทัด (ถ้าบรรทัดซ้อนกันให้เพิ่มเลขนี้, ถ้าห่างไปให้ลดเลขนี้)
    gap = 60        

    # --- เริ่มเขียนข้อมูลทีละบรรทัด ---
    
    # บรรทัดที่ 1: วันเวลา
    draw.text((start_x, start_y), data["{{DATE}}"], font=font_text, fill=text_color)
    
    # บรรทัดที่ 2: สถานที่
    draw.text((start_x, start_y + gap), data["{{LOCATION}}"], font=font_text, fill=text_color)
    
    # บรรทัดที่ 3: ประเภท
    draw.text((start_x, start_y + gap*2), data["{{TYPE}}"], font=font_text, fill=text_color)
    
    # บรรทัดที่ 4: ผู้ควบคุม
    draw.text((start_x, start_y + gap*3), data["{{COMMANDER}}"], font=font_text, fill=text_color)
    
    # บรรทัดที่ 5: ระดับความเสี่ยง (ตรงนี้ที่ Error แก้ให้แล้วครับ)
    draw.text((start_x, start_y + gap*4), data["{{RISK}}"], font=font_text, fill=text_color)
    
    # บรรทัดที่ 6: ตรวจยานพาหนะ
    draw.text((start_x, start_y + gap*5), data["{{VEHICLE}}"], font=font_text, fill=text_color)
    
    # บรรทัดที่ 7: ผู้ติดต่อ (ข้าม 1 จังหวะเพื่อให้ตรงช่อง)
    draw.text((start_x, start_y + gap*6), data["{{COORD_NAME}}"], font=font_text, fill=text_color)
    
    # บรรทัดที่ 8: พิกัด
    draw.text((start_x, start_y + gap*7), data["{{GPS}}"], font=font_text, fill=text_color)
    
    # บรรทัดที่ 9: เส้นทาง/สถานการณ์
    draw.text((start_x, start_y + gap*8), data["{{SITUATION}}"], font=font_text, fill=text_color)
    # 4. แปะรูปภาพ 4 รูป (ฝั่งซ้าย)
    # พิกัดกรอบรูป (สมมติ)
    # รูป 1 (ซ้ายบน) | รูป 2 (ขวาบน)
    # รูป 3 (ซ้ายล่าง) | รูป 4 (ขวาล่าง)
    
    # กำหนดขนาดรูปที่ต้องการย่อ (เช่น 350x250 pixel)
    target_size = (350, 250) 
    
    # พิกัดมุมซ้ายบนของแต่ละรูป (X, Y)
    positions = [
        (50, 250),   # รูปที่ 1
        (420, 250),  # รูปที่ 2
        (50, 520),   # รูปที่ 3
        (420, 520)   # รูปที่ 4
    ]

    for i, img_file in enumerate(images):
        if i < 4:
            # เปิดไฟล์รูป
            photo = Image.open(img_file)
            # ย่อรูปให้พอดี
            photo = photo.resize(target_size)
            # แปะลงบนพื้นหลัง
            base_img.paste(photo, positions[i])

    return base_img

# --- ส่วนรับข้อมูลจาก User ---
st.markdown("---")
st.subheader("ส่วนหัวกระดาษ")
header_month = st.text_input("หัวข้อเรื่อง", value="สายตรวจท่าเรือประจำวัน ประจำเดือน พ.ย.68")
st.markdown("---")

col1, col2 = st.columns(2)

with col1:
    st.subheader("ข้อมูลการตรวจ")
    date_time = st.text_input("วัน เวลา ออกตรวจ", value=f"{datetime.now().strftime('%d/%m/%Y')} เวลาประมาณ 09.30 น.")
    location = st.text_input("สถานที่", value="ท่าเทียบเรือ ต.เจ๊ะเห อ.ตากใบ จ.นราธิวาส")
    type_port = st.text_input("ประเภท", value="ท่าเรือ")
    commander = st.text_input("ผู้ควบคุม", value="พ.ต.ท.จิรายุทธ์ แก้วด้วง สว.ส.รน.4 กก.7 บก.รน.")
    risk_level = st.selectbox("ระดับความเสี่ยง", ["สีเขียว", "สีเหลือง", "สีแดง"])

with col2:
    st.subheader("รายละเอียดเพิ่มเติม")
    vehicle = st.text_input("ตรวจยานพาหนะ", value="เรือข้ามฟาก, แพขนานยนต์ ไทย-มาเลเซีย")
    coordinator = st.text_input("ผู้ติดต่อประสานงาน", value="- ผู้ดูแล")
    coordinates = st.text_input("พิกัด", value="6.235873N, 102.08970241E")
    situation = st.text_area("เส้นทาง/สถานการณ์", value="ทางบก / เหตุการณ์ทั่วไปปกติ")

st.markdown("---")
st.subheader("ภาพประกอบ (เลือก 4 รูป)")
uploaded_files = st.file_uploader("เลือกรูปภาพ", type=['jpg', 'png'], accept_multiple_files=True)

# Preview รูป
if uploaded_files:
    use_files = uploaded_files[:4]
    cols = st.columns(4)
    for i, img_file in enumerate(use_files):
        with cols[i]:
            st.image(img_file, caption=f"รูปที่ {i+1}", use_container_width=True)

# --- ส่วนปุ่ม Download ---
st.markdown("### 📥 เลือกรูปแบบไฟล์ที่ต้องการดาวน์โหลด")
d_col1, d_col2, d_col3 = st.columns(3)

# รวบรวมข้อมูล
data_dict = {
    "{{HEADER_MONTH}}": header_month,
    "{{DATE}}": date_time,
    "{{LOCATION}}": location,
    "{{TYPE}}": type_port,
    "{{COMMANDER}}": commander,
    "{{RISK}}": risk_level,
    "{{VEHICLE}}": vehicle,
    "{{COORD_NAME}}": coordinator,
    "{{GPS}}": coordinates,
    "{{SITUATION}}": situation
}

# 1. ปุ่ม PowerPoint
with d_col1:
    if st.button("Download PowerPoint (.pptx)"):
        if not uploaded_files:
            st.warning("⚠️ กรุณาใส่รูปภาพก่อนครับ")
        else:
            prs = Presentation(required_files["template"])
            slide = prs.slides[0]
            # แก้ไขข้อความ
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for key, val in data_dict.items():
                        replace_text_ppt(shape, key, val)
            # ใส่รูป
            images_to_insert = uploaded_files[:4]
            img_index = 0
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 18:
                    if img_index < len(images_to_insert):
                        shape.insert_picture(images_to_insert[img_index])
                        img_index += 1
            
            out_ppt = io.BytesIO()
            prs.save(out_ppt)
            out_ppt.seek(0)
            st.download_button("คลิกเพื่อโหลด PPTX", out_ppt, f"Report_{header_month}.pptx")

# 2. ปุ่ม Image (JPG)
with d_col2:
    if st.button("Download Image (.jpg)"):
        if not uploaded_files:
            st.warning("⚠️ กรุณาใส่รูปภาพก่อนครับ")
        else:
            final_img = generate_image_report(data_dict, uploaded_files[:4], required_files["background"], required_files["font"])
            if final_img:
                out_jpg = io.BytesIO()
                final_img.save(out_jpg, format="JPEG", quality=95)
                out_jpg.seek(0)
                st.download_button("คลิกเพื่อโหลด JPG", out_jpg, f"Report_{header_month}.jpg", mime="image/jpeg")

# 3. ปุ่ม PDF
with d_col3:
    if st.button("Download PDF (.pdf)"):
        if not uploaded_files:
            st.warning("⚠️ กรุณาใส่รูปภาพก่อนครับ")
        else:
            final_img = generate_image_report(data_dict, uploaded_files[:4], required_files["background"], required_files["font"])
            if final_img:
                out_pdf = io.BytesIO()
                # แปลงภาพเป็น PDF
                final_img.save(out_pdf, format="PDF", resolution=100.0)
                out_pdf.seek(0)
                st.download_button("คลิกเพื่อโหลด PDF", out_pdf, f"Report_{header_month}.pdf", mime="application/pdf")









